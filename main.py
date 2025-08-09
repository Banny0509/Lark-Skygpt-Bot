# Lark SkyGPT - FastAPI webhook with multi-modal, @mention gating in groups, P2P auto-reply,
# and daily 08:00 summary for previous day's group chats (Asia/Taipei by default).
# --- OPTIMIZED VERSION ---
import os
import json
import logging
import re
import base64
from io import BytesIO
from datetime import datetime, timezone, timedelta
from typing import Optional, Dict, List, Tuple, Any
from collections import defaultdict
import asyncio
from sys import exit as sys_exit # 用於啟動失敗時退出

import httpx
from fastapi import FastAPI, Request, HTTPException
import aiosqlite  # [OPTIMIZED] 使用 aiosqlite 進行異步資料庫操作

# Optional parsers for office/PDF and images
try:
    import pdfplumber  # type: ignore
except Exception:
    pdfplumber = None
try:
    import openpyxl  # type: ignore
except Exception:
    openpyxl = None
try:
    import docx  # type: ignore  # python-docx
except Exception:
    docx = None
try:
    import pptx  # type: ignore  # python-pptx
except Exception:
    pptx = None
try:
    from PIL import Image  # type: ignore
except Exception:
    Image = None

logging.basicConfig(level=logging.INFO,
                    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s")
logger = logging.getLogger("lark-skygpt")

# Load .env locally when not on Railway (or other prod)
if "RAILWAY_ENVIRONMENT" not in os.environ:
    try:
        from dotenv import load_dotenv  # type: ignore
        load_dotenv()
        logger.info("Loaded .env for local development")
    except Exception:
        pass

app = FastAPI()

# --- Environment ---
LARK_APP_ID = os.getenv("APP_ID")              # Bot APP_ID (string like cli_aaxxx)
LARK_APP_SECRET = os.getenv("APP_SECRET")      # Bot APP_SECRET
VERIFICATION_TOKEN = os.getenv("VERIFICATION_TOKEN")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o")
TIMEZONE_OFFSET = int(os.getenv("TIMEZONE_OFFSET", "8"))  # Asia/Taipei UTC+8 by default
DATABASE_FILE = os.getenv("DATABASE_FILE", "lark_chat_history.db")  # 資料庫檔名

# [IMAGE/FORWARD] Cache for the most recent media (image or file) in each chat.  
# Keyed by chat_id. Each value is a dict storing type ("image" or "file"),
# key (image_key or file_key), optional file_name, message_id, and timestamp (datetime).
LAST_MEDIA_CACHE: Dict[str, Dict[str, Any]] = {}
# TTL for cached media in seconds. If user asks a question within this window, we attach the last media.
MEDIA_CACHE_TTL = 180  # 3 minutes

# Name used to mention the bot in group chats. This should match the display name of your bot
# in Lark (e.g., "Skygpt"). The name is case-insensitive and whitespace trimmed.
BOT_NAME = os.getenv("BOT_NAME", "").strip()

# [OPTIMIZED] 快速失敗：如果缺少關鍵配置，直接中止程式
if not all([LARK_APP_ID, LARK_APP_SECRET, VERIFICATION_TOKEN, OPENAI_API_KEY]):
    logger.error("One or more required env vars are missing. Application cannot start.")
    sys_exit("Fatal: Missing required environment variables.")

# [OPTIMIZED] 全域共用的 HTTP 客戶端，由 FastAPI 生命週期管理
http_client: Optional[httpx.AsyncClient] = None

# --- [OPTIMIZED] 異步資料庫函式 (Async Database Functions) ---
async def init_db():
    """初始化資料庫和資料表"""
    async with aiosqlite.connect(DATABASE_FILE) as con:
        await con.execute("""
            CREATE TABLE IF NOT EXISTS chat_history (
                message_id TEXT PRIMARY KEY,
                chat_id TEXT NOT NULL,
                chat_type TEXT NOT NULL,
                timestamp TEXT NOT NULL,
                text TEXT,
                summarized INTEGER DEFAULT 0
            )
        """)
        await con.commit()
    logger.info("Database initialized successfully: %s", DATABASE_FILE)

async def log_message_to_db(message_id: str, chat_id: str, chat_type: str, timestamp: datetime, text: str):
    """將訊息非阻塞地紀錄到資料庫"""
    async with aiosqlite.connect(DATABASE_FILE) as con:
        try:
            await con.execute(
                "INSERT INTO chat_history (message_id, chat_id, chat_type, timestamp, text) VALUES (?, ?, ?, ?, ?)",
                (message_id, chat_id, chat_type, timestamp.isoformat(), text)
            )
            await con.commit()
        except aiosqlite.IntegrityError:
            logger.warning("Message with ID %s already exists.", message_id)
        except Exception as e:
            logger.exception("Failed to log message to DB: %s", e)

async def get_chats_for_summary(start_date: datetime, end_date: datetime) -> Dict[str, List[Dict]]:
    """從資料庫非阻塞地獲取需要摘要的群組聊天紀錄"""
    chats = defaultdict(list)
    async with aiosqlite.connect(DATABASE_FILE) as con:
        con.row_factory = aiosqlite.Row
        async with con.execute(
            "SELECT chat_id, timestamp, text FROM chat_history "
            "WHERE chat_type = 'group' AND summarized = 0 AND timestamp BETWEEN ? AND ? "
            "ORDER BY timestamp ASC",
            (start_date.isoformat(), end_date.isoformat())
        ) as cursor:
            async for row in cursor:
                chats[row['chat_id']].append({
                    "timestamp": datetime.fromisoformat(row['timestamp']),
                    "text": row['text']
                })
    return chats

async def mark_chat_messages_as_summarized(chat_id: str, start_date: datetime, end_date: datetime):
    """[OPTIMIZED] 將指定聊天室的已摘要訊息在資料庫中標記"""
    async with aiosqlite.connect(DATABASE_FILE) as con:
        await con.execute(
            "UPDATE chat_history SET summarized = 1 WHERE chat_id = ? AND timestamp BETWEEN ? AND ?",
            (chat_id, start_date.isoformat(), end_date.isoformat())
        )
        await con.commit()
    logger.info("Marked messages for chat_id %s from %s to %s as summarized.", chat_id, start_date, end_date)

# --- 時間工具與健康檢查 ---
def now_local() -> datetime:
    return datetime.now(timezone(timedelta(hours=TIMEZONE_OFFSET)))

def format_datetime_info(dt: Optional[datetime] = None) -> str:
    dt = dt or now_local()
    tz_sign = "+" if TIMEZONE_OFFSET >= 0 else "-"
    tz_abs = abs(TIMEZONE_OFFSET)
    weekday_map = ["星期一","星期二","星期三","星期四","星期五","星期六","星期日"]
    weekday = weekday_map[dt.isoweekday()-1]
    return f"{dt.strftime('%Y-%m-%d %H:%M:%S')} (Asia/Taipei, UTC{tz_sign}{tz_abs:02d}:00, {weekday})"

@app.get("/health")
async def health():
    return {"status": "ok", "now": format_datetime_info()}

# --- Lark helper functions (using shared client) ---
async def get_lark_token() -> str:
    if not http_client:
        raise RuntimeError("HTTP client not initialized.")
    resp = await http_client.post(
        "https://open.larksuite.com/open-apis/auth/v3/tenant_access_token/internal",
        json={"app_id": LARK_APP_ID, "app_secret": LARK_APP_SECRET},
    )
    resp.raise_for_status()
    data = resp.json()
    token = data.get("tenant_access_token")
    if not token:
        raise RuntimeError(f"Failed to get tenant_access_token: {data}")
    return token

async def send_message_to_lark(chat_id: str, text: str) -> None:
    if not http_client:
        raise RuntimeError("HTTP client not initialized.")
    token = await get_lark_token()
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json; charset=utf-8",
    }
    payload = {"receive_id": chat_id, "msg_type": "text", "content": json.dumps({"text": text}, ensure_ascii=False)}
    r = await http_client.post(
        "https://open.larksuite.com/open-apis/im/v1/messages?receive_id_type=chat_id",
        headers=headers, json=payload
    )
    try:
        r.raise_for_status()
    except Exception:
        logger.error("send_message_to_lark failed: %s", r.text)

async def download_message_resource(message_id: str, file_key: str) -> bytes:
    if not http_client:
        raise RuntimeError("HTTP client not initialized.")
    token = await get_lark_token()
    url = f"https://open.larksuite.com/open-apis/im/v1/messages/{message_id}/resources/{file_key}"
    headers = {"Authorization": f"Bearer {token}"}
    r = await http_client.get(url, headers=headers)
    # On certain tenants, images cannot be downloaded via the generic resource endpoint. In that case,
    # fall back to the dedicated image endpoint using the file_key as image_key.
    if r.status_code == 404 or r.status_code == 403:
        # Try image endpoint (note: some tenants require image_type param; default works for origin)
        fallback_url = f"https://open.larksuite.com/open-apis/im/v1/images/{file_key}"
        r2 = await http_client.get(fallback_url, headers=headers)
        r2.raise_for_status()
        return r2.content
    r.raise_for_status()
    return r.content

# --- Parsing helpers ---
def _strip_mentions(text: str) -> str:
    return re.sub(r"<at[^>]*>.*?</at>", "", text or "").strip()

def extract_text_from_file(file_bytes: bytes, file_name: str) -> str:
    # ... (此函式內容不變) ...
    name_lower = (file_name or "").lower()
    try:
        if (name_lower.endswith(".xlsx") or name_lower.endswith(".xls")) and openpyxl is not None:
            wb = openpyxl.load_workbook(filename=BytesIO(file_bytes), data_only=True, read_only=True)
            sheet = wb[wb.sheetnames[0]]
            lines = []
            for row in sheet.iter_rows(values_only=True):
                line = "\t".join("" if c is None else str(c) for c in row)
                lines.append(line)
            return "\n".join(lines)
        if name_lower.endswith(".docx") and docx is not None:
            d = docx.Document(BytesIO(file_bytes))
            return "\n".join(p.text for p in d.paragraphs)
        if (name_lower.endswith(".pptx") or name_lower.endswith(".ppt")) and pptx is not None:
            prs = pptx.Presentation(BytesIO(file_bytes))
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        out.append(shape.text)
            return "\n".join(out)
        if name_lower.endswith(".pdf") and pdfplumber is not None:
            texts = []
            with pdfplumber.open(BytesIO(file_bytes)) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        texts.append(t)
            return "\n".join(texts) if texts else "[空白 PDF 或無法抽取文字]"
    except Exception as e:
        logger.exception("extract_text_from_file error for %s: %s", file_name, e)
        return f"[解析 {file_name} 發生錯誤: {e}]"
    return f"[不支援的檔案或缺少解析庫: {file_name}]"


def build_image_content(b64_png: str, user_hint: str = "請閱讀此圖片並回覆重點或回答問題。") -> List[dict]:
    return [
        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64_png}", "detail": "low"}},
        {"type": "text", "text": user_hint},
    ]

# --- OpenAI (using shared client) ---
async def call_openai_api(messages: List[dict], model: Optional[str] = None) -> dict:
    if not http_client:
        raise RuntimeError("HTTP client not initialized.")
    headers = {"Authorization": f"Bearer {OPENAI_API_KEY}", "Content-Type": "application/json"}
    data = {"model": model or OPENAI_MODEL, "messages": messages}
    r = await http_client.post("https://api.openai.com/v1/chat/completions", headers=headers, json=data)
    r.raise_for_status()
    return r.json()

def system_datetime_message() -> dict:
    return {
        "role": "system",
        "content": (
            "You are SkyGPT. "
            f"Current local datetime is: {format_datetime_info()}. "
            "Timezone is Asia/Taipei (UTC+08:00). "
            "When users ask about the current date/time/day, ALWAYS answer using this clock. "
            "Do not say you cannot access time."
        ),
    }

# --- Webhook ---
@app.post("/webhook")
async def webhook(request: Request):
    # ... (此函式內容不變) ...
    body = await request.body()
    try:
        payload = json.loads(body.decode("utf-8"))
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON")

    if "challenge" in payload:
        return {"challenge": payload["challenge"]}

    header = payload.get("header", {})
    if header.get("token") != VERIFICATION_TOKEN:
        raise HTTPException(status_code=403, detail="Invalid token")

    if header.get("event_type") == "im.message.receive_v1":
        # 使用 asyncio.create_task 讓 webhook 立即回傳，避免 Lark 端超時
        asyncio.create_task(handle_message_receive(payload.get("event", {})))

    return {"code": 0}

# --- 更穩健的 @ 提及檢測 ---
def _is_bot_mentioned_in_group(message: dict) -> bool:
    # ... (此函式內容不變) ...
    if message.get("chat_type") != "group":
        return False

    mentions = message.get("mentions", []) or []
    for m in mentions:
        # Collect all possible id values from the mention (user_id, open_id, app_id, etc.)
        possible_ids = set()
        id_info = m.get("id", {})
        if isinstance(id_info, dict):
            for v in id_info.values():
                if isinstance(v, str):
                    possible_ids.add(v)
        app_id_top = m.get("app_id")
        if isinstance(app_id_top, str):
            possible_ids.add(app_id_top)
        # Check if this mention matches the bot's known IDs
        if LARK_APP_ID and LARK_APP_ID in possible_ids:
            return True
        # Additionally check if the mention name matches the configured BOT_NAME
        mention_name = (m.get("name") or "").strip()
        if BOT_NAME and mention_name.lower() == BOT_NAME.lower():
            return True

    raw = message.get("content") or "{}"
    try:
        c = json.loads(raw)
        txt = c.get("text", "")
    except Exception:
        txt = ""
    # Check <at> tags for id attributes matching the bot's APP_ID
    for attr in ("id", "user_id", "open_id", "app_id"):
        for m in re.finditer(fr'{attr}="([^\"]+)"', txt):
            if LARK_APP_ID and m.group(1) == LARK_APP_ID:
                return True
    # Finally, check for a plain-text @ mention of the bot's name
    if BOT_NAME:
        # Use word boundary after the name to avoid partial matches
        pattern = re.compile(r'@' + re.escape(BOT_NAME) + r'\b', re.IGNORECASE)
        if pattern.search(txt):
            return True
    return False

# --- 主訊息處理 ---
async def handle_message_receive(event: dict):
    message = event.get("message", {})
    if not message:
        return

    logger.info("Received message payload: %s", json.dumps(message))

    chat_type = message.get("chat_type")
    chat_id = message.get("chat_id")
    message_id = message.get("message_id")

    create_time_ms = message.get("create_time")
    ts_local = datetime.fromtimestamp(int(create_time_ms) / 1000, tz=timezone.utc).astimezone(timezone(timedelta(hours=TIMEZONE_OFFSET))) if create_time_ms else now_local()

    msg_type = message.get("message_type")
    content = json.loads(message.get("content") or "{}")

    # 記錄到 SQLite
    summary_text = ""
    if msg_type == "text":
        summary_text = _strip_mentions(content.get("text", ""))
    elif msg_type == "image":
        summary_text = "[圖片]"
    elif msg_type == "file":
        summary_text = f"[檔案: {content.get('file_name') or ''}]"
    if all([message_id, chat_id, chat_type, summary_text]):
        await log_message_to_db(message_id, chat_id, chat_type, ts_local, summary_text)

    # 更新媒體快取：記錄該聊天室最近的圖片或檔案（3分鐘內有效）
    try:
        if msg_type == "image":
            image_key = content.get("image_key")
            if image_key:
                LAST_MEDIA_CACHE[chat_id] = {
                    "type": "image",
                    "key": image_key,
                    "message_id": message_id,
                    "timestamp": ts_local,
                }
        elif msg_type == "file":
            file_key = content.get("file_key")
            file_name = content.get("file_name") or ""
            if file_key:
                LAST_MEDIA_CACHE[chat_id] = {
                    "type": "file",
                    "key": file_key,
                    "name": file_name,
                    "message_id": message_id,
                    "timestamp": ts_local,
                }
    except Exception as e:
        logger.exception("Error updating media cache: %s", e)

    # 觸發條件
    is_trigger = (chat_type == "p2p") or (chat_type == "group" and _is_bot_mentioned_in_group(message))
    if not is_trigger:
        logger.info("Ignored message (no @mention in group or not P2P).")
        return

    # 處理時間問題
    user_text_for_dt = _strip_mentions(content.get("text", "")).lower() if msg_type == "text" else ""
    dt_patterns = [r"今天|日期|幾號|現在幾點|現在時間|現在|時間|幾點|星期幾|禮拜幾", r"\b(date|today|time|current)\b"]
    if user_text_for_dt and any(re.search(p, user_text_for_dt) for p in dt_patterns):
        await send_message_to_lark(chat_id, f"現在時間：{format_datetime_info()}")
        return

    # 其餘分流處理
    try:
        reply = ""
        if msg_type == "image":
            image_key = content.get("image_key")
            if not image_key:
                await send_message_to_lark(chat_id, "收到圖片但缺少 image_key。")
                return
            file_bytes = await download_message_resource(message_id, image_key)
            b64 = base64.b64encode(file_bytes).decode("utf-8")
            mm = [system_datetime_message(), {"role": "user", "content": build_image_content(b64)}]
            resp = await call_openai_api(mm)
            reply = resp["choices"][0]["message"]["content"].strip()

        elif msg_type == "file":
            file_key, file_name = content.get("file_key"), content.get("file_name") or "附件"
            if not file_key:
                await send_message_to_lark(chat_id, "收到檔案但缺少 file_key。")
                return
            file_bytes = await download_message_resource(message_id, file_key)
            extracted = extract_text_from_file(file_bytes, file_name)
            prompt = f"請閱讀以下檔案內容，使用繁體中文摘要重點：\n\n{extracted[:15000]}"
            mm = [system_datetime_message(), {"role": "user", "content": prompt}]
            resp = await call_openai_api(mm)
            reply = resp["choices"][0]["message"]["content"].strip()

        elif msg_type == "text":
            user_text = _strip_mentions(content.get("text", ""))
            if not user_text:
                await send_message_to_lark(chat_id, "（空白訊息）")
                return
            # 檢查最近快取的圖片/檔案，若在有效時間內則附帶處理
            last_media = LAST_MEDIA_CACHE.get(chat_id)
            mm = [system_datetime_message()]
            attached = False
            if last_media:
                try:
                    age = (ts_local - last_media.get("timestamp", ts_local)).total_seconds()
                    if age <= MEDIA_CACHE_TTL:
                        if last_media.get("type") == "image":
                            # 讀取並附上上一張圖片內容
                            file_bytes = await download_message_resource(last_media["message_id"], last_media["key"])
                            b64 = base64.b64encode(file_bytes).decode("utf-8")
                            mm.append({"role": "user", "content": build_image_content(b64)})
                            attached = True
                        elif last_media.get("type") == "file":
                            # 讀取並附上上一個檔案內容
                            file_bytes = await download_message_resource(last_media["message_id"], last_media["key"])
                            extracted = extract_text_from_file(file_bytes, last_media.get("name") or "")
                            # 附檔案內容前幾萬字，避免過長
                            mm.append({"role": "user", "content": f"以下是近期附件內容：\n{extracted[:15000]}"})
                            attached = True
                except Exception as e:
                    logger.exception("Failed to fetch cached media: %s", e)
                    # 移除快取以避免下次重複失敗
                    LAST_MEDIA_CACHE.pop(chat_id, None)
            # 最後加入使用者的問題文字
            mm.append({"role": "user", "content": user_text})
            resp = await call_openai_api(mm)
            reply = resp["choices"][0]["message"]["content"].strip()
        else:
            await send_message_to_lark(chat_id, "這個訊息類型暫不支援，請傳文字、圖片或檔案。")
            return
        
        if reply:
            await send_message_to_lark(chat_id, reply)

    except Exception as e:
        logger.exception("Error handling message: %s", e)
        await send_message_to_lark(chat_id, f"抱歉，處理訊息時發生錯誤：{e}")


# --- [OPTIMIZED] 每日 08:00 摘要 (邏輯更穩健) ---
async def daily_summary_scheduler():
    while True:
        try:
            now_dt = now_local()
            next_run = now_dt.replace(hour=8, minute=0, second=0, microsecond=0)
            if now_dt >= next_run:
                next_run += timedelta(days=1)
            wait_seconds = (next_run - now_dt).total_seconds()
            logger.info("Scheduler: Next summary run at %s (in %.2f hours)", next_run, wait_seconds / 3600)
            await asyncio.sleep(wait_seconds)

            run_time = now_local()
            day_start = (run_time - timedelta(days=1)).replace(hour=0, minute=0, second=0, microsecond=0)
            day_end = (run_time - timedelta(days=1)).replace(hour=23, minute=59, second=59, microsecond=999999)

            chats_to_summarize = await get_chats_for_summary(day_start, day_end)
            if not chats_to_summarize:
                logger.info("No new group chats to summarize for %s.", day_start.strftime('%Y-%m-%d'))
                continue

            for chat_id, messages in chats_to_summarize.items():
                if not messages:
                    continue
                
                content_lines = [f"{m['timestamp'].strftime('%H:%M')}: {m['text']}" for m in messages]
                combined = "\n".join(content_lines)
                prompt = "請將以下群組聊天記錄整理成摘要（繁體中文，條列式，含關鍵決策、待辦、未決問題）：\n\n" + combined[:15000]
                
                try:
                    resp = await call_openai_api([system_datetime_message(), {"role": "user", "content": prompt}])
                    summary = resp["choices"][0]["message"]["content"].strip()
                    await send_message_to_lark(chat_id, f"昨日聊天摘要 ({day_start.strftime('%Y-%m-%d')}):\n{summary}")
                    
                    # 成功發送摘要後，才標記這個聊天室的訊息為已處理
                    await mark_chat_messages_as_summarized(chat_id, day_start, day_end)

                except Exception as e:
                    logger.exception("Summary generation/sending failed for chat_id %s. Will retry next cycle. Error: %s", chat_id, e)
                    # 失敗後不標記，以便下次重試

        except Exception as e:
            logger.exception("daily_summary_scheduler error: %s", e)
            await asyncio.sleep(60) # 發生重大錯誤時，等待1分鐘後重試

# --- FastAPI 生命週期事件 ---
@app.on_event("startup")
async def startup_event():
    """應用啟動時執行的事件"""
    global http_client
    http_client = httpx.AsyncClient(timeout=60) # 建立共用的 client
    await init_db() # 初始化資料庫
    asyncio.create_task(daily_summary_scheduler()) # 啟動背景排程

@app.on_event("shutdown")
async def shutdown_event():
    """應用關閉時執行的事件"""
    if http_client:
        await http_client.aclose() # 優雅地關閉 client

# --- Web server entry (local run) ---
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=int(os.getenv("PORT", "8080")))
